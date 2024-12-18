import json
import os
import requests
from pptx import Presentation as PPTXPresentation
from uuid import uuid4

API_KEY = os.getenv("SLIDE_SPEAK_API_KEY")

class Presentation:
    def __init__(self, file=None, url=None):
        if not file and not url:
            raise ValueError("You must provide either a file path or a URL.")

        self.file_path = file if file else self._download_presentation(url)
        self.slides_data = self._extract_presentation_data()

    def _download_presentation(self, url):
        print("Downloading presentation...")
        response = requests.get(url)
        if response.status_code != 200:
            raise Exception(f"Failed to download file: {response.status_code}")
        temp_file = f"presentation_{uuid4().hex}.pptx"
        with open(temp_file, "wb") as f:
            f.write(response.content)
        print(f"Download complete! Saved to {temp_file}")
        return temp_file

    def _extract_presentation_data(self):
        presentation = PPTXPresentation(self.file_path)
        slides_data = []
        for i, slide in enumerate(presentation.slides):
            slide_data = {"slide_number": i + 1, "shapes": []}
            for shape in slide.shapes:
                text = "<No text>"
                name = shape.name if hasattr(shape, "name") else "Unnamed Shape"
                if shape.has_text_frame and shape.text.strip():
                    text = shape.text.strip()
                slide_data["shapes"].append({"name": name, "content": text})
            slides_data.append(slide_data)
        return slides_data

    def save_json(self, output_file=None):
        if not output_file:
            output_file = self.file_path.replace(".pptx", ".json")
        with open(output_file, "w", encoding="utf-8") as json_file:
            json.dump({"presentation": self.slides_data}, json_file, indent=2)
        print(f"✅ JSON content saved to '{output_file}'.")

    def get_slides(self):
        return self.slides_data

    @classmethod
    def edit_presentation(cls, original_file, edits_json):
        api_url = "https://api.slidespeak.co/api/v1/presentation/edit"
        headers = {"X-API-key": API_KEY}
        print("Sending edits to the presentation...")
        print({"config": json.dumps(edits_json)})

        try:
            with open(original_file, "rb") as f:
                response = requests.post(
                    api_url,
                    headers=headers,
                    files={"pptx_file": f},
                    data={"config": json.dumps(edits_json)}
                )
        except FileNotFoundError:
            print(f"❌ File not found: {original_file}")
            return None

        if response.status_code == 200:
            response_json = response.json()
            new_presentation_url = response_json.get("url")
            if new_presentation_url:
                print(f"✅ Edits applied successfully! New presentation available at: {new_presentation_url}")
                return cls(url=new_presentation_url)
            else:
                raise Exception("❌ API response does not contain a URL.")
        else:
            raise Exception(f"❌ Failed to edit presentation: {response.status_code}, {response.text}")
